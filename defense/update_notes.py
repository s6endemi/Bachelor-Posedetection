# -*- coding: utf-8 -*-
"""Synct speech_script.md in die Notizenfelder von defense_slides.pptx.

Fasst NUR die Notizen der Folien 1-23 an — die Folien selbst (auch manuelle
PowerPoint-Aenderungen) und die Backup-Notizen (B1-B9) bleiben unberuehrt.

Nutzung: nach jeder Aenderung an speech_script.md ausfuehren:
    .venv/Scripts/python.exe defense/update_notes.py
"""
import re
import os
from pptx import Presentation

DEF = os.path.dirname(os.path.abspath(__file__))
SCRIPT = os.path.join(DEF, "speech_script.md")
PPTX = os.path.join(DEF, "defense_slides.pptx")

md = open(SCRIPT, encoding="utf-8").read()

# Abschnitte: "## Folie N — ..." bis zum naechsten "## " oder "# "
blocks = re.findall(r"## Folie (\d+)[^\n]*\n(.*?)(?=\n## |\n# |\Z)", md, re.S)

prs = Presentation(PPTX)
slides = list(prs.slides)

updated = 0
for num, body in blocks:
    i = int(num) - 1
    if i >= len(slides):
        continue
    text = body.replace("\n---\n", "\n").replace("---", "").strip()
    slides[i].notes_slide.notes_text_frame.text = text
    updated += 1

prs.save(PPTX)
print(f"OK: {updated} Folien-Notizen aus speech_script.md uebernommen -> {os.path.basename(PPTX)}")
print("Folien selbst und Backup-Notizen (B1-B9) unveraendert.")
