# -*- coding: utf-8 -*-
"""Generiert defense_slides.pptx aus der Storyline (defense/slides_storyline.md).

Stufe-2-Werkzeug: Inhalte hier aendern -> neu ausfuehren.
Nach Beginn des PowerPoint-Feinschliffs (Stufe 3) NICHT mehr neu generieren,
sonst gehen manuelle Aenderungen verloren (Output wird ueberschrieben).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))  # bachelor/
FIG = os.path.join(ROOT, "thesis", "fig")
DEF = os.path.join(ROOT, "defense")

# ---- Design ----------------------------------------------------------------
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x77, 0x77, 0x77)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
MP_GREEN = RGBColor(0x2E, 0x7D, 0x32)
MN_RED = RGBColor(0xC6, 0x28, 0x28)
YL_BLUE = RGBColor(0x15, 0x65, 0xC0)
ACCENT = RGBColor(0x1F, 0x38, 0x64)
FONT = "Calibri"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

page_counter = {"n": 0}


def slide(page_label=None):
    s = prs.slides.add_slide(BLANK)
    page_counter["n"] += 1
    label = page_label if page_label else str(page_counter["n"])
    tb = s.shapes.add_textbox(SW - Inches(1.0), SH - Inches(0.42), Inches(0.85), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.RIGHT
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.color.rgb = GRAY
    r.font.name = FONT
    return s


def title(s, text, color=DARK, size=27, kicker=None):
    y0 = Inches(0.28)
    if kicker:
        kb = s.shapes.add_textbox(Inches(0.5), Inches(0.16), SW - Inches(1.0), Inches(0.35))
        kp = kb.text_frame.paragraphs[0]
        kp.text = kicker
        kr = kp.runs[0]
        kr.font.size = Pt(13)
        kr.font.bold = True
        kr.font.color.rgb = GRAY
        kr.font.name = FONT
        y0 = Inches(0.52)
    tb = s.shapes.add_textbox(Inches(0.5), y0, SW - Inches(1.0), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    r.font.name = FONT
    return tb


def text(s, x, y, w, h, content, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
         line_spacing=1.15):
    """content: str oder Liste von (text, dict)-Zeilen bzw. Strings."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = content if isinstance(content, list) else [content]
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            t, opts = line
        else:
            t, opts = line, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t if t else " "
        p.alignment = opts.get("align", align)
        p.line_spacing = line_spacing
        p.space_after = Pt(opts.get("space_after", 6))
        r = p.runs[0]
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.color.rgb = opts.get("color", color)
        r.font.name = FONT
    return tb


def box(s, x, y, w, h, fill=LIGHT, line_color=None, txt=None, size=16, bold=False,
        color=DARK, align=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.08
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    if txt is not None:
        lines = txt if isinstance(txt, list) else [txt]
        for i, line in enumerate(lines):
            if isinstance(line, tuple):
                t, opts = line
            else:
                t, opts = line, {}
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = t if t else " "
            p.alignment = opts.get("align", align)
            r = p.runs[0]
            r.font.size = Pt(opts.get("size", size))
            r.font.bold = opts.get("bold", bold)
            r.font.color.rgb = opts.get("color", color)
            r.font.name = FONT
    return sh


def arrow(s, x1, y1, x2, y2):
    conn = s.shapes.add_connector(2, int(x1), int(y1), int(x2), int(y2))  # straight
    conn.line.color.rgb = GRAY
    conn.line.width = Pt(2.25)
    le = conn.line._get_or_add_ln()
    import copy
    from lxml import etree
    tail = etree.SubElement(le, ('{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd'))
    tail.set('type', 'triangle')
    return conn


def pic_fit(s, path, x, y, max_w, max_h, center_in_box=True):
    im = Image.open(path)
    iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    px = x + (max_w - w) / 2 if center_in_box else x
    py = y + (max_h - h) / 2 if center_in_box else y
    return s.shapes.add_picture(path, int(px), int(py), int(w), int(h))


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


MODEL_LINE = [("MediaPipe", MP_GREEN), ("MoveNet", MN_RED), ("YOLOv8", YL_BLUE)]

# =============================================================================
# 1 — Titel
# =============================================================================
s = slide()
pic_fit(s, os.path.join(DEF, "Logo_Bonn.png"), Inches(11.4), Inches(0.3), Inches(1.5), Inches(1.5), False)
text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.8), [
    ("A Holistic Evaluation of 2D Mobile Pose Estimation Models", {"size": 32, "bold": True, "color": ACCENT}),
    ("for Home-Based Rehabilitation", {"size": 32, "bold": True, "color": ACCENT}),
])
text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(2.4), [
    ("Bachelor Thesis Defense — Eren Demir", {"size": 20, "bold": True}),
    ("Supervisor: Lars Doorenbos   ·   Second examiner: Prof. Dr. Jürgen Gall", {"size": 16, "color": GRAY}),
    ("University of Bonn — Institute of Computer Science, Computer Vision Group", {"size": 16, "color": GRAY}),
    ("July 2026", {"size": 16, "color": GRAY}),
])
notes(s, "Begruessung. Ein Satz: I evaluated three mobile pose estimation models for home-based "
         "rehabilitation across six quality dimensions on a motion-capture-validated benchmark.")

# =============================================================================
# 2 — Motivation
# =============================================================================
s = slide()
title(s, "Home rehabilitation fails without feedback")
text(s, Inches(0.5), Inches(2.0), Inches(6.2), Inches(4.6), [
    ("Adherence to home exercise programs is poor (Jack et al. 2010)", {}),
    ("Non-adherence prolongs recovery and raises costs", {}),
    ("A single RGB camera + pose estimation = markerless exercise monitoring — no markers, no suits, no lab", {}),
    ("Open question: which mobile model can carry such a system?", {"bold": True, "color": ACCENT}),
], size=18)
bx, by, bw, bh, gap = Inches(7.2), Inches(2.0), Inches(5.4), Inches(0.95), Inches(0.55)
box(s, bx, by, bw, bh, LIGHT, None, [("Patient exercises at home", {"bold": True})])
arrow(s, bx + bw / 2, by + bh, bx + bw / 2, by + bh + gap)
box(s, bx, by + bh + gap, bw, bh, LIGHT, None, [("Smartphone camera + pose estimation model", {"bold": True})])
arrow(s, bx + bw / 2, by + bh + gap + bh, bx + bw / 2, by + 2 * (bh + gap))
box(s, bx, by + 2 * (bh + gap), bw, bh, LIGHT, None, [("Exercise feedback & progress monitoring", {"bold": True})])
notes(s, "OPENER (auswendig): 'Physiotherapy works when exercises are done correctly and "
         "regularly - but most of it happens at home, unsupervised, and adherence is poor. "
         "A smartphone camera with a pose estimation model could close that gap. My thesis "
         "asks: which of today's mobile models can you actually trust in a living room? And "
         "the answer turned out to be more interesting than a ranking.'")

# =============================================================================
# 3 — Constraints
# =============================================================================
s = slide()
title(s, "Four deployment constraints that standard benchmarks do not cover")
cw, ch = Inches(5.9), Inches(1.85)
cells = [
    ("Camera positioning", "Patients rarely face the camera frontally — oblique and lateral views, self-occlusion"),
    ("Multi-person interference", "Family members or a therapist enter the frame — the target subject must be retained"),
    ("Temporal stability", "Real-time video feedback needs stable predictions — jitter ruins usability"),
    ("Resource constraints", "No GPU at home — usable throughput on CPU is required"),
]
for i, (h, b) in enumerate(cells):
    x = Inches(0.5) + (i % 2) * (cw + Inches(0.45))
    y = Inches(1.5) + (i // 2) * (ch + Inches(0.35))
    box(s, x, y, cw, ch, LIGHT, None, [
        (h, {"bold": True, "size": 18, "color": ACCENT}),
        (b, {"size": 15}),
    ], align=PP_ALIGN.LEFT)
text(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.9),
     [("COCO and MPII do not combine these conditions  →  the comparison has to be made on rehabilitation-style data",
       {"bold": True, "size": 18, "color": ACCENT})])
notes(s, "Kern-Symmetrie ansagen: each of these four constraints becomes an evaluation dimension later. "
         "Standard-Benchmarks sind wertvoll fuer Architekturvergleich, decken aber dieses Setting nicht ab.")

# =============================================================================
# 4 — Research Questions
# =============================================================================
s = slide()
title(s, "Five research questions")
rqs = [
    ("RQ1", "How accurate and complete are the models on rehabilitation data with MoCap ground truth?"),
    ("RQ2", "How temporally stable are their predictions?"),
    ("RQ3", "How robust are they under suboptimal conditions — body orientation and multi-person interference?"),
    ("RQ4", "How far do rankings transfer to a standard benchmark (COCO)?"),
    ("RQ5", "What practical trade-offs and recommendations follow for deployment?"),
]
y = Inches(1.55)
for tag, q in rqs:
    box(s, Inches(0.5), y, Inches(1.0), Inches(0.85), ACCENT, None, [(tag, {"bold": True, "color": RGBColor(255, 255, 255), "size": 17})])
    text(s, Inches(1.7), y + Inches(0.06), Inches(11.1), Inches(0.9), q, size=17)
    y += Inches(1.02)
notes(s, "Merkformel: how good, how stable, how robust, does it transfer, what to deploy. "
         "Jede RQ wird in der Discussion explizit beantwortet.")

# =============================================================================
# 5 — Paradigmen
# =============================================================================
s = slide()
title(s, "Three models, three architectural paradigms")
pw, ph = Inches(4.0), Inches(4.15)
paradigms = [
    ("Top-down", "MediaPipe Pose (BlazePose)", MP_GREEN,
     ["detect person → crop → estimate pose", "Most precise per person (full-res crop)",
      "Cost grows with person count", "Depends on the upstream detector", "33 native landmarks"]),
    ("Bottom-up", "MoveNet (MultiPose)", MN_RED,
     ["one full-image pass → keypoints grouped via person centers", "Cost independent of person count",
      "Grouping is the weak spot", "17 COCO keypoints + boxes"]),
    ("One-stage", "YOLOv8-Pose", YL_BLUE,
     ["object detector with a pose head", "Box + 17 keypoints in one shot",
      "Best speed integration", "No high-resolution per-person zoom"]),
]
for i, (para, model, col, desc) in enumerate(paradigms):
    x = Inches(0.5) + i * (pw + Inches(0.4))
    box(s, x, Inches(1.5), pw, Inches(0.62), col, None, [(para, {"bold": True, "color": RGBColor(255, 255, 255), "size": 18})])
    lines = [(model, {"bold": True, "size": 16, "color": col})] + [(d, {"size": 14}) for d in desc]
    box(s, x, Inches(2.12), pw, ph - Inches(0.62), LIGHT, col, lines, align=PP_ALIGN.LEFT)
text(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.7),
     [("Selection criteria: CPU real-time capability · mature deployment ecosystem · paradigm coverage    (+ 6 family variants)",
       {"size": 15, "color": GRAY})])
notes(s, "GALL-FOLIE. Auswahl kriteriengeleitet, nicht willkuerlich. Ein Vertreter pro Paradigma -> "
         "Unterschiede architektonisch interpretierbar (Kap. 6). Varianten: MP Lite/Heavy, MoveNet SP Lightning/Thunder, YOLO s/m.")

# =============================================================================
# 6 — REHAB24-6
# =============================================================================
s = slide()
title(s, "REHAB24-6: a rehabilitation-oriented proxy benchmark")
pic_fit(s, os.path.join(ROOT, "validation_skeleton.png"), Inches(0.4), Inches(1.45), Inches(7.3), Inches(4.6))
text(s, Inches(7.95), Inches(1.5), Inches(5.0), Inches(4.6), [
    ("10 healthy adults (6 m / 4 f, age 25–50)", {}),
    ("6 physiotherapist-guided exercises, 1,072 repetitions", {}),
    ("2 synchronized RGB views (c17 frontal, c18 lateral)", {}),
    ("Ground truth: 41-marker MoCap → 26-joint skeleton → projected to 2D", {}),
    ("63 recordings → 126 camera-view sequences", {}),
    ("Green = 12 evaluated joints (shared skeleton)", {"color": MP_GREEN, "size": 15}),
], size=16)
text(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.7),
     [("A proxy testbed for rehabilitation-style movement analysis — not clinical validation",
       {"bold": True, "size": 17, "color": ACCENT})])
notes(s, "Proxy-Framing von Anfang an (Schutzschild). Bild: echtes GT-Skelett auf echtem Frame — "
         "und im Hintergrund die Therapeutin (Multi-Person-Vorschau).")

# =============================================================================
# 7 — Protokoll
# =============================================================================
s = slide()
title(s, "A fixed evaluation protocol isolates the models")
steps = [
    ("Video", "(every 3rd frame, 10 Hz)"),
    ("Frame-independent inference", "(no tracking, no smoothing)"),
    ("Person selection", "(largest instance)"),
    ("12-joint mapping", "(33 / 17 / 26 → 12)"),
    ("Joint filter", "(confidence ≥ 0.5)"),
    ("NMPJPE + displacement", ""),
]
bw2, bh2 = Inches(1.95), Inches(1.5)
x = Inches(0.35)
y = Inches(1.7)
for i, (st_a, st_b) in enumerate(steps):
    lines = [(st_a, {"size": 12.5, "bold": i == 1})]
    if st_b:
        lines.append((st_b, {"size": 11.5, "color": GRAY}))
    box(s, x, y, bw2, bh2, LIGHT if i not in (1,) else RGBColor(0xE3, 0xEB, 0xF7), None, lines)
    if i < len(steps) - 1:
        arrow(s, x + bw2, y + bh2 / 2, x + bw2 + Inches(0.18), y + bh2 / 2)
    x += bw2 + Inches(0.18)
text(s, Inches(0.5), Inches(3.7), Inches(12.3), Inches(1.1), [
    ("NMPJPE  =  mean joint error / ground-truth torso length  (×100%)", {"bold": True, "size": 19, "color": ACCENT}),
    ("scale-free — comparable across subjects, camera distances, resolutions", {"size": 15, "color": GRAY}),
])
text(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(1.1), [
    ("Why frame-independent?  We evaluate models, not pipelines.", {"size": 17, "bold": True}),
], size=16)
notes(s, "Versteckte Komplexitaet zeigen. Frame-independent: YOLO hat kein natives temporales Processing -> "
         "jeder andere Modus waere unfair; temporal processing kann uniform on top. Werte = upper bound der "
         "Roh-Outputs. Muendlich ergaenzen: zwei Regimes - central tendency (valid subset) vs. failure "
         "behavior (Outlier >100% Torso separat analysiert).")

# =============================================================================
# 8 — Statistik als Problem→Lösung-Kette
# =============================================================================
s = slide()
title(s, "Frames lie about sample size — honest statistics")
stat_rows = [
    ("367,000 frames — but frames of one person are not independent",
     "aggregate to 10 subject clusters"),
    ("Only ten data points per model",
     "exact permutation test: all 1024 sign flips, no distribution assumptions"),
    ("Three pairwise comparisons",
     "Holm correction controls false alarms"),
    ("Significant ≠ large",
     "bootstrap confidence intervals for effect size"),
]
y = Inches(1.65)
for prob, sol in stat_rows:
    box(s, Inches(0.5), y, Inches(6.1), Inches(1.05), LIGHT, None,
        [(prob, {"size": 15, "bold": True})], align=PP_ALIGN.LEFT)
    arrow(s, Inches(6.7), y + Inches(0.52), Inches(7.05), y + Inches(0.52))
    box(s, Inches(7.15), y, Inches(5.65), Inches(1.05), RGBColor(0xE3, 0xEB, 0xF7), None,
        [(sol, {"size": 15})], align=PP_ALIGN.LEFT)
    y += Inches(1.22)
text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
     [("+ auxiliary COCO val2017 comparison — different protocol, kept descriptive", {"size": 14, "color": GRAY})])
notes(s, "Warum-Kette sprechen: jede Wahl folgt aus einem Problem. Median pro Sequenz -> Mittel pro Person -> "
         "10 Zahlen pro Modell. Details Backup B1. "
         "SIGNPOST danach: 'That's the instrument. Six dimensions, nine configurations, one protocol. "
         "Here is what we found.'")

# =============================================================================
# 9 — Accuracy
# =============================================================================
s = slide()
title(s, "MoveNet is the most accurate — significant in all 10 subject clusters",
      kicker="DIMENSION 1 OF 6 — BODY-JOINT ACCURACY")
pic_fit(s, os.path.join(FIG, "fig_accuracy_boxplot.png"), Inches(0.35), Inches(1.4), Inches(7.6), Inches(5.4))
text(s, Inches(8.2), Inches(1.7), Inches(4.7), Inches(4.9), [
    ("Mean NMPJPE (valid subset):", {"bold": True, "size": 17}),
    ("MoveNet      10.52%", {"color": MN_RED, "bold": True, "size": 19}),
    ("MediaPipe   12.53%", {"color": MP_GREEN, "size": 19}),
    ("YOLOv8       12.77%", {"color": YL_BLUE, "size": 19}),
    ("", {}),
    ("All pairwise differences survive Holm correction", {"size": 15}),
    ("MoveNet wins in all 10 clusters — p = 2/1024, the exact-test minimum", {"size": 15, "bold": True}),
    ("First on all 6 exercises", {"size": 15}),
], size=16)
notes(s, "Boxplot: jeder Punkt = eine Person (Cluster-NMPJPE). MoveNets ganze Verteilung liegt unter den anderen. "
         "Cluster-Differenzen: MP-MoveNet +1.01pp [0.62,1.42]; MoveNet-YOLO -1.50pp [-1.71,-1.27]; MP-YOLO -0.49pp p=.0352.")

# =============================================================================
# 10 — Accuracy per Joint (Dimension 1b, Hip-Befund)
# =============================================================================
s = slide()
title(s, "The errors are not uniform — MediaPipe's hips stand out",
      kicker="DIMENSION 1 OF 6 — ACCURACY BY BODY REGION")
pj_rows_main = [
    ("Region", "MediaPipe", "MoveNet", "YOLOv8"),
    ("Shoulders", "8.8%", "10.3%", "10.2%"),
    ("Elbows", "11.2%", "10.0%", "13.1%"),
    ("Wrists", "12.0%", "11.8%", "16.3%"),
    ("Hips", "16.7%", "11.3%", "14.6%"),
    ("Knees", "9.9%", "7.8%", "9.9%"),
    ("Ankles", "14.6%", "11.4%", "12.6%"),
]
y = Inches(1.75)
for i, (a, b, c, d) in enumerate(pj_rows_main):
    bold = i == 0
    hip = a == "Hips"
    fill = RGBColor(0xE3, 0xEB, 0xF7) if i == 0 else (RGBColor(0xFB, 0xE9, 0xE7) if hip else (LIGHT if i % 2 else RGBColor(0xFF, 0xFF, 0xFF)))
    box(s, Inches(0.5), y, Inches(2.4), Inches(0.62), fill, None, [(a, {"size": 14, "bold": bold or hip})], align=PP_ALIGN.LEFT)
    box(s, Inches(2.9), y, Inches(1.9), Inches(0.62), fill, None, [(b, {"size": 14, "bold": bold or hip, "color": MP_GREEN if not bold else DARK})])
    box(s, Inches(4.8), y, Inches(1.9), Inches(0.62), fill, None, [(c, {"size": 14, "bold": bold or hip, "color": MN_RED if not bold else DARK})])
    box(s, Inches(6.7), y, Inches(1.9), Inches(0.62), fill, None, [(d, {"size": 14, "bold": bold, "color": YL_BLUE if not bold else DARK})])
    y += Inches(0.62)
text(s, Inches(8.95), Inches(1.85), Inches(4.0), Inches(4.6), [
    ("MediaPipe wins only the shoulders", {"size": 15, "color": MP_GREEN, "bold": True}),
    ("MoveNet leads everywhere else", {"size": 15, "color": MN_RED}),
    ("", {}),
    ("Hips: 16.7% vs. 11.3% — a concentrated problem, not general weakness", {"size": 15, "bold": True}),
    ("Excluding hips shrinks the MediaPipe–MoveNet gap by ~40% (2.0 → 1.2 pp)", {"size": 14}),
    ("Consistent with a landmark-definition mismatch — indirect evidence", {"size": 14, "color": GRAY}),
], size=15)
notes(s, "Hip-Befund: konzentriert, nicht uniform. 3 Hueft-Konventionen: COCO-Annotationen (MoveNet/YOLO-Training), "
         "BlazePose-eigenes Schema, MoCap-Gelenkzentrum (GT). Ohne Hips: Gap 2.01->1.21pp. Als Hypothese "
         "kennzeichnen (indirect). Mapping selbst ist visuell validiert.")

# =============================================================================
# 10 — Viewpoint
# =============================================================================
s = slide()
title(s, "All models degrade on lateral views — MoveNet the least",
      kicker="DIMENSION 2 OF 6 — VIEWPOINT SENSITIVITY")
pic_fit(s, os.path.join(FIG, "fig_rotation_line.png"), Inches(5.6), Inches(1.45), Inches(7.4), Inches(5.3))
text(s, Inches(0.5), Inches(1.8), Inches(5.0), Inches(4.5), [
    ("Frontal (0–20°) → lateral (60–90°):", {"bold": True, "size": 17}),
    ("MoveNet      9.75% → 11.49%   (+17.8%)", {"color": MN_RED, "size": 17, "bold": True}),
    ("MediaPipe   11.28% → 14.38%   (+27.5%)", {"color": MP_GREEN, "size": 17}),
    ("YOLOv8      11.14% → 14.66%   (+31.7%)", {"color": YL_BLUE, "size": 17}),
    ("", {}),
    ("Orientation from 3D MoCap shoulders — no circularity", {"size": 14, "color": GRAY}),
    ("Fixed camera geometry, bimodal viewpoints — a viewpoint comparison, not a dense rotation study", {"size": 14, "color": GRAY}),
], size=16)
notes(s, "Winkel aus GT, nicht aus Predictions (Zirkularitaet). Globaler 65-Grad-Offset (+-5), 10-Grad-Bins > Unsicherheit; "
         "Vergleich zwischen Modellen invariant gegen Offset-Fehler. Details Backup B3.")

# =============================================================================
# 11 — Stability
# =============================================================================
s = slide()
title(s, "MoveNet is also the most stable — the whole family leads",
      kicker="DIMENSION 3 OF 6 — TEMPORAL STABILITY")
pic_fit(s, os.path.join(FIG, "fig_prediction_displacement_boxplot.png"), Inches(0.35), Inches(1.4), Inches(7.6), Inches(5.4))
text(s, Inches(8.2), Inches(1.7), Inches(4.7), Inches(4.9), [
    ("Frame-to-frame displacement (% torso):", {"bold": True, "size": 16}),
    ("MoveNet      3.81%", {"color": MN_RED, "bold": True, "size": 19}),
    ("YOLOv8       5.22%", {"color": YL_BLUE, "size": 19}),
    ("MediaPipe   6.03%", {"color": MP_GREEN, "size": 19}),
    ("", {}),
    ("Across all nine variants, the MoveNet family takes the top three (3.31 / 3.61 / 3.81)", {"size": 15, "bold": True}),
    ("Raw, unsmoothed outputs — true motion is identical across models, so differences isolate prediction noise", {"size": 14, "color": GRAY}),
], size=16)
notes(s, "Proxy in 1 Satz: displacement = prediction vs. its own prediction one frame later; GT nur als Lineal "
         "(Torso-Normalisierung). Inhaerente Stabilitaet, kein Filter-Verdienst.")

# =============================================================================
# 12 — Speed (Phase A, vierter MoveNet-Sieg)
# =============================================================================
s = slide()
title(s, "MoveNet leads on speed — family accuracy gains cost real latency",
      kicker="DIMENSION 4 OF 6 — CPU INFERENCE SPEED")
pic_fit(s, os.path.join(FIG, "fig_speed_accuracy_scatter.png"), Inches(0.35), Inches(1.5), Inches(8.2), Inches(5.3))
text(s, Inches(8.8), Inches(1.8), Inches(4.2), Inches(4.9), [
    ("CPU throughput (main configs):", {"bold": True, "size": 16}),
    ("MoveNet MultiPose   27.7 FPS", {"color": MN_RED, "size": 16, "bold": True}),
    ("YOLOv8 Nano            18.9 FPS", {"color": YL_BLUE, "size": 16}),
    ("MediaPipe Full          14.7 FPS", {"color": MP_GREEN, "size": 16}),
    ("", {}),
    ("Fastest: MoveNet SP Lightning — 100.6 FPS", {"size": 15}),
    ("Accuracy leaders pay: MediaPipe Heavy 6.0 · YOLOv8 Medium 4.6 FPS", {"size": 15}),
], size=15)
notes(s, "Dashed frontier = non-dominated configs. Native Stacks (Tasks API / TFLite / PyTorch) — deployment-"
         "realistisch. SIGNPOST danach (Pause!): 'Four dimensions in, MoveNet sweeps — accuracy, viewpoint, "
         "stability, speed. A simple story so far. The remaining two dimensions break it.'")

# =============================================================================
# 12 — Completeness
# =============================================================================
s = slide()
title(s, "YOLOv8 delivers the complete skeleton most often",
      kicker="DIMENSION 5 OF 6 — DETECTION COMPLETENESS")
bars = [("YOLOv8", 79.1, YL_BLUE), ("MediaPipe", 55.1, MP_GREEN), ("MoveNet", 38.2, MN_RED)]
by = Inches(2.0)
for name, val, col in bars:
    text(s, Inches(0.6), by, Inches(1.9), Inches(0.5), [(name, {"bold": True, "size": 17, "color": col})])
    bar_w = Emu(int(Inches(9.0) * (val / 100.0)))
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.6), by, bar_w, Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb = col; sh.line.fill.background(); sh.shadow.inherit = False
    text(s, Inches(2.7) + bar_w, by, Inches(1.2), Inches(0.5), [(f"{val}%", {"bold": True, "size": 17})])
    by += Inches(1.0)
text(s, Inches(0.6), Inches(5.3), Inches(12.0), Inches(1.2), [
    ("Share of frames with all 12 evaluated joints after confidence filtering", {"size": 15, "color": GRAY}),
    ("First dimension with a different winner — the profiles start to separate", {"bold": True, "size": 18, "color": ACCENT}),
], size=16)
notes(s, "Mean valid joints: 11.53 / 10.72 / 10.31 (YOLO/MP/MoveNet). Uebergang: bisher gewann MoveNet alles — ab hier nicht mehr.")

# =============================================================================
# 13 — Zwei Regime
# =============================================================================
s = slide()
title(s, "A second person in frame is harmless — until it is the coach",
      kicker="DIMENSION 6 OF 6 — MULTI-PERSON BEHAVIOR")
box(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(3.4), LIGHT, None, [
    ("Dataset-wide multi-person subset", {"bold": True, "size": 18, "color": ACCENT}),
    ("(metadata severity ≥ 2, ~10,000 frames per model)", {"size": 13, "color": GRAY}),
    ("", {}),
    ("Accuracy barely changes:", {"size": 15}),
    ("MediaPipe +1.2%   ·   MoveNet −5.0%   ·   YOLOv8 −4.4%", {"size": 15, "bold": True}),
    ("(composition effects — not causal)", {"size": 13, "color": GRAY}),
], align=PP_ALIGN.LEFT)
box(s, Inches(6.9), Inches(1.7), Inches(6.0), Inches(3.4), RGBColor(0xFB, 0xE9, 0xE7), None, [
    ("Coach extreme case", {"bold": True, "size": 18, "color": MN_RED}),
    ("(5 sequences with a persistently visible therapist)", {"size": 13, "color": GRAY}),
    ("", {}),
    ("Catastrophic person switches appear —", {"size": 15}),
    ("outlier rates 9.1% / 13.8% / 13.9%", {"size": 15, "bold": True}),
    ("(MediaPipe / MoveNet / YOLOv8)", {"size": 13, "color": GRAY}),
], align=PP_ALIGN.LEFT)
text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.0),
     [("Multi-person behavior must be decomposed into two regimes — averaging them away hides the real deployment risk",
       {"bold": True, "size": 18, "color": ACCENT})])
notes(s, "Contribution 5. Links: blosse Anwesenheit ist harmlos. Rechts: anhaltend prominente Person = echtes Risiko.")

# =============================================================================
# 14 — Coach case
# =============================================================================
s = slide()
title(s, "One frame, three behaviors")
pic_fit(s, os.path.join(FIG, "coach_detection_comparison.png"), Inches(0.4), Inches(1.35), Inches(12.5), Inches(4.9))
text(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.8), [
    ("2+ persons detected: MediaPipe 23% · MoveNet 16% · YOLOv8 62%      |      coach outlier rate: 9.1% · 13.8% · 13.9%",
     {"bold": True, "size": 16})])
notes(s, "PM_119-c17 Frame 105, illustrativ. MediaPipe: nur Patientin. MoveNet: detektiert beide, largest-area waehlt "
         "den Therapeuten -> SELECTION-stage failure, nicht Detection. YOLO: 3 Kandidaten, waehlt richtig. Aggregate in der Zeile unten.")

# =============================================================================
# 15 — Detection count != robustness
# =============================================================================
s = slide()
title(s, "The obvious explanation fails — in three steps")
rows = [
    ("MediaPipe vs. YOLOv8", "MediaPipe reports a 2nd person in 23% of frames, YOLOv8 in 62% — more detections, more failures", "consistent with a detection-level explanation", LIGHT),
    ("MediaPipe vs. MoveNet", "MoveNet detects the 2nd person LESS often (16% < 23%) yet fails MORE (13.8% > 9.1%)", "breaks the explanation — mechanism beyond detection frequency", RGBColor(0xFB, 0xE9, 0xE7)),
    ("MoveNet vs. YOLOv8", "4× difference in 2nd-person detections (16% vs. 62%)", "nearly identical failure rates (13.8% vs. 13.9%)", LIGHT),
]
y = Inches(1.55)
for name, fact, verdict, fill in rows:
    box(s, Inches(0.5), y, Inches(3.1), Inches(1.35), ACCENT, None, [(name, {"bold": True, "color": RGBColor(255, 255, 255), "size": 15})])
    box(s, Inches(3.75), y, Inches(9.1), Inches(1.35), fill, None, [
        (fact, {"size": 15, "bold": True}), (verdict, {"size": 14, "color": GRAY})], align=PP_ALIGN.LEFT)
    y += Inches(1.55)
text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.9),
     [("Detection count alone does not explain coach robustness — the residual mechanism sits at the selection stage and is left open",
       {"size": 15, "bold": True, "color": ACCENT})])
notes(s, "Der staerkste interpretative Befund — bewusst eine NEGATIVE Konklusion (was es NICHT ist). "
         "Panel (c) der Coach-Figur zeigt den Selection-Fehler konkret.")

# =============================================================================
# 18 — Coach pro Video (PM_010: Ehrlichkeits-Beat)
# =============================================================================
s = slide()
title(s, "One video breaks all three — no selection strategy is perfect")
coach_rows_main = [
    ("Coach sequence", "MediaPipe", "MoveNet", "YOLOv8"),
    ("PM_010-c17", "81.9%", "70.7%", "69.9%"),
    ("PM_011-c17", "38.0%", "59.7%", "66.2%"),
    ("PM_108-c17", "17.2%", "46.2%", "51.1%"),
    ("PM_119-c17", "24.1%", "74.0%", "73.0%"),
    ("PM_121-c17", "30.9%", "60.3%", "75.3%"),
]
y = Inches(1.8)
for i, (a, b, c, d) in enumerate(coach_rows_main):
    bold = i == 0
    pm10 = a == "PM_010-c17"
    fill = RGBColor(0xE3, 0xEB, 0xF7) if i == 0 else (RGBColor(0xFB, 0xE9, 0xE7) if pm10 else (LIGHT if i % 2 else RGBColor(0xFF, 0xFF, 0xFF)))
    box(s, Inches(0.7), y, Inches(2.9), Inches(0.68), fill, None, [(a, {"size": 15, "bold": bold or pm10})], align=PP_ALIGN.LEFT)
    box(s, Inches(3.6), y, Inches(2.2), Inches(0.68), fill, None, [(b, {"size": 15, "bold": bold or pm10, "color": MP_GREEN if not bold else DARK})])
    box(s, Inches(5.8), y, Inches(2.2), Inches(0.68), fill, None, [(c, {"size": 15, "bold": bold or pm10, "color": MN_RED if not bold else DARK})])
    box(s, Inches(8.0), y, Inches(2.2), Inches(0.68), fill, None, [(d, {"size": 15, "bold": bold or pm10, "color": YL_BLUE if not bold else DARK})])
    y += Inches(0.68)
text(s, Inches(10.45), Inches(1.9), Inches(2.5), Inches(3.5), [
    ("Mean coach NMPJPE per video (before outlier cut)", {"size": 13, "color": GRAY}),
], size=13)
text(s, Inches(0.7), Inches(6.15), Inches(12.0), Inches(1.1), [
    ("In four of five videos, MediaPipe stays below 40% — but PM_010 breaks all three models. MediaPipe is the most robust — not immune.", {"bold": True, "size": 16, "color": ACCENT}),
], size=15)
notes(s, "Ehrlichkeits-Beat: PM_010 = Gegenprobe, bricht ALLE drei (MP 81.9!). MediaPipe nur most robust, "
         "nicht immun. Zahlen VOR Outlier-Cut (failure regime). n=5, deskriptiv — benannte Limitation.")

# =============================================================================
# 17 — Failure modes
# =============================================================================
s = slide()
title(s, "Each architecture fails in its own way")
pic_fit(s, os.path.join(DEF, "fig_failure_mode_taxonomy.png"), Inches(0.35), Inches(1.4), Inches(7.4), Inches(5.2))
text(s, Inches(8.0), Inches(1.6), Inches(4.9), Inches(5.0), [
    ("MediaPipe → Keypoint-Displacement", {"color": MP_GREEN, "bold": True, "size": 16}),
    ("confident misprediction — 69.6% of 1,583 failures", {"size": 14, "color": GRAY}),
    ("MoveNet → Confidence-Collapse", {"color": MN_RED, "bold": True, "size": 16}),
    ("graceful degradation — 51.4% of 1,875", {"size": 14, "color": GRAY}),
    ("YOLOv8 → Missing-Detection", {"color": YL_BLUE, "bold": True, "size": 16}),
    ("no detection at all — 63.8% of 5,610 (≈20× MediaPipe's no-detection rate)", {"size": 14, "color": GRAY}),
    ("", {}),
    ("Consistent with each paradigm's design logic — observational, one family per paradigm", {"size": 14}),
], size=15)
notes(s, "Hypothesen-Level klar halten: observational, nicht mechanistisch — Architektur konfundiert mit Trainingsdaten/"
         "Kalibrierung/Landmark-Definitionen. Robust fuer Cutoff 4-8.")

# =============================================================================
# 18 — Cross-dataset
# =============================================================================
s = slide()
title(s, "Rankings only partially transfer to COCO")
pic_fit(s, os.path.join(FIG, "fig_cross_dataset_slopegraph.png"), Inches(0.35), Inches(1.4), Inches(8.2), Inches(5.4))
text(s, Inches(8.8), Inches(1.7), Inches(4.2), Inches(4.9), [
    ("REHAB #1 MoveNet → COCO #4", {"size": 15, "bold": True}),
    ("COCO #1: YOLOv8 Medium", {"size": 15}),
    ("MediaPipe drops sharply (5 → 8)", {"size": 15}),
    ("Main models: MoveNet first on both; MediaPipe & YOLOv8 swap", {"size": 15}),
    ("", {}),
    ("Different protocol → limited transferability, not benchmark invalidity", {"size": 14, "color": GRAY}),
], size=15)
notes(s, "Shift bleibt auch bei Single-Person-COCO. Protokoll: torso-center matching zur GT, nur sichtbare Joints, deskriptiv. Backup B6.")

# =============================================================================
# 19 — Drei Profile
# =============================================================================
s = slide()
title(s, "Three architectural profiles — not one ranking")
profs = [
    ("MediaPipe — filtered top-down", MP_GREEN,
     ["Lowest coach failure rate (9.1%)", "33 native landmarks, simple API", "",
      "— but slower (14.7 FPS) and weakest hips (16.7%)"]),
    ("MoveNet — low-exposure bottom-up", MN_RED,
     ["Best accuracy (10.5%)", "Lowest displacement (3.8%)", "Highest main-model FPS (27.7)", "",
      "— but lowest completeness (38%) and not the safest in coach scenes"]),
    ("YOLOv8 — high-sensitivity one-stage", YL_BLUE,
     ["Most complete detections (79%)", "Most stable across datasets", "",
      "— but detects other people most often (62%), weakest viewpoint robustness"]),
]
for i, (name, col, desc) in enumerate(profs):
    x = Inches(0.5) + i * (Inches(4.0) + Inches(0.4))
    box(s, x, Inches(1.55), Inches(4.0), Inches(0.62), col, None, [(name, {"bold": True, "color": RGBColor(255, 255, 255), "size": 15})])
    box(s, x, Inches(2.17), Inches(4.0), Inches(3.6), LIGHT, col, [(d, {"size": 14.5}) for d in desc], align=PP_ALIGN.LEFT)
text(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.8),
     [("Which model is right depends on which dimension your application cannot afford to lose",
       {"bold": True, "size": 18, "color": ACCENT})])
notes(s, "BRUECKEN-SATZ zum Einstieg: 'So back to the original question — which model for home rehabilitation?' "
         "DIE Kernfolie, Spiegel von Folie 5 (gleiche Spaltenreihenfolge wie die Paradigmen!). Kein universeller "
         "Sieger — drei Deployment-Profile. Decision-relevant answer, keine Diplomatie.")

# =============================================================================
# 20 — Recommendations
# =============================================================================
s = slide()
title(s, "Recommendations for CPU-based deployment")
recs = [
    ("Default (body joints + CPU)", "MoveNet MultiPose — accuracy + stability + throughput", MN_RED),
    ("Completeness is critical", "YOLOv8 — full 12-joint skeleton most often", YL_BLUE),
    ("Coach robustness / richer landmarks", "MediaPipe — lowest catastrophic coach rate, 33 landmarks", MP_GREEN),
    ("Maximum throughput", "MoveNet SP Lightning — 100 FPS, if accuracy loss is acceptable", MN_RED),
]
y = Inches(1.55)
for headline, rec, col in recs:
    box(s, Inches(0.5), y, Inches(4.6), Inches(0.95), LIGHT, None, [(headline, {"bold": True, "size": 15})], align=PP_ALIGN.LEFT)
    box(s, Inches(5.25), y, Inches(7.6), Inches(0.95), RGBColor(0xFF, 0xFF, 0xFF), col, [(rec, {"size": 15})], align=PP_ALIGN.LEFT)
    y += Inches(1.12)
box(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.85), RGBColor(0xE3, 0xEB, 0xF7), None,
    [("Deployment guidance matters as much as model choice: frontal camera positioning · interference monitoring · quality warnings",
      {"size": 15, "bold": True})])
notes(s, "RQ5-Antwort. Der Bonus-Satz unten ist der Abbinder aus 6.7.")

# =============================================================================
# 21 — Limitations
# =============================================================================
s = slide()
title(s, "Limitations are scope decisions — each with a next step")
lims = [
    ("Healthy volunteers, 10 subject clusters", "→ validation on real patient cohorts (next step #1)"),
    ("Frame-independent raw outputs (upper bound)", "→ uniform temporal post-processing applied to all models"),
    ("COCO comparison under a different protocol", "→ standardized cross-dataset evaluation protocol"),
    ("2D joint localization only", "→ 3D joint angles and movement quality"),
]
y = Inches(1.65)
for lim, nxt in lims:
    box(s, Inches(0.5), y, Inches(6.4), Inches(1.0), LIGHT, None, [(lim, {"size": 15, "bold": True})], align=PP_ALIGN.LEFT)
    box(s, Inches(7.05), y, Inches(5.75), Inches(1.0), RGBColor(0xE8, 0xF5, 0xE9), None, [(nxt, {"size": 15})], align=PP_ALIGN.LEFT)
    y += Inches(1.18)
notes(s, "Nicht defensiv: bewusste Scope-Grenzen. Nie eine Schwaeche ohne den naechsten Schritt. "
         "Weitere Limitationen (Coach n=5, Displacement-Proxy, CPU-only) muendlich falls gefragt.")

# =============================================================================
# 22 — Conclusion
# =============================================================================
s = slide()
title(s, "Conclusion")
box(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.55), LIGHT, None, [
    ("1 — MoveNet MultiPose is the strongest CPU default for body-joint rehabilitation analysis", {"bold": True, "size": 18}),
    ("lowest subject-cluster error, lowest displacement, highest main-model throughput", {"size": 14, "color": GRAY}),
], align=PP_ALIGN.LEFT)
box(s, Inches(0.5), Inches(3.35), Inches(12.3), Inches(1.85), LIGHT, None, [
    ("2 — The three models form three architectural profiles, not one ranking", {"bold": True, "size": 18}),
    ("detection count alone does not explain coach-scene robustness; each architecture carries a distinct failure signature", {"size": 14, "color": GRAY}),
], align=PP_ALIGN.LEFT)
text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.9),
     [("The most important next step toward clinical use: validation on real patient data",
       {"bold": True, "size": 18, "color": ACCENT})])
notes(s, "Conclusion-Kapitel kondensiert. Ruhig und klar enden — dann Fragen.")

# =============================================================================
# 23 — Thank you
# =============================================================================
s = slide()
pic_fit(s, os.path.join(DEF, "Logo_Bonn.png"), Inches(11.6), Inches(0.3), Inches(1.3), Inches(1.3), False)
text(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(1.2),
     [("Thank you — questions?", {"size": 36, "bold": True, "color": ACCENT})])
text(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(1.2), [
    ("A Holistic Evaluation of 2D Mobile Pose Estimation Models for Home-Based Rehabilitation", {"size": 16, "color": GRAY}),
    ("Eren Demir · University of Bonn · July 2026", {"size": 14, "color": GRAY}),
])
notes(s, "CLOSER (auswendig, auf der Conclusion-Folie sprechen): 'So — can you trust a mobile pose estimator "
         "in a living room? Yes. But which one depends on what your application cannot afford to lose. MoveNet "
         "MultiPose is the default; the three architectures fail in three different ways, and knowing those "
         "ways is as valuable as any ranking. The next step toward the clinic is validation on real patients. "
         "Thank you.' — Backups: B1 Statistik, B2 NMPJPE, B3 Rotation, B4 Mapping, B5 Varianten, B6 COCO, "
         "B7 Coach/Video, B8 Hips, B9 Architektur.")

# =============================================================================
# Backups
# =============================================================================
def backup(label):
    return slide(page_label=label)

# B1 Statistik
s = backup("B1")
title(s, "Backup: statistical analysis in detail", ACCENT)
text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.3), [
    ("Aggregation: median NMPJPE per camera-view sequence (robust to residual outliers) → mean of medians per subject cluster (person_id) → 10 values per model", {"size": 16}),
    ("Paired sign-flip permutation test: under H₀ each cluster difference's sign is a coin flip → enumerate all 2¹⁰ = 1024 sign assignments → exact two-sided p-value; minimum attainable p = 2/1024 ≈ 0.002", {"size": 16}),
    ("Holm correction: sort p-values ascending, test against α/k, α/(k−1), … — same family-wise error control as Bonferroni, uniformly more power", {"size": 16}),
    ("Percentile bootstrap: resample the 10 cluster differences with replacement, 20,000×, take the 2.5/97.5 percentiles → effect size + uncertainty", {"size": 16}),
    ("Results: MP−MoveNet +1.01pp [0.62, 1.42], p_Holm=.0059 · MoveNet−YOLO −1.50pp [−1.71, −1.27], p_Holm=.0059 · MP−YOLO −0.49pp [−0.85, −0.12], p_Holm=.0352", {"size": 15, "bold": True}),
    ("Video-level sensitivity analysis: same paired test, Monte-Carlo sign sampling (2¹²⁶ not enumerable) — agrees with cluster level", {"size": 16}),
], size=16)

# B2 NMPJPE
s = backup("B2")
title(s, "Backup: NMPJPE definition and metric alternatives", ACCENT)
text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.3), [
    ("NMPJPE = (1/N) Σ ‖pᵢ − gᵢ‖₂ / L_torso × 100%,   N = 12 evaluated joints", {"size": 19, "bold": True}),
    ("L_torso = distance between GT shoulder midpoint and GT hip midpoint (per frame)", {"size": 16}),
    ("Joints below confidence 0.5 are excluded per joint; frames keep their remaining valid joints", {"size": 16}),
    ("", {}),
    ("Why not the alternatives:", {"bold": True, "size": 16}),
    ("PCK — collapses error size into hit/miss around a threshold", {"size": 15}),
    ("OKS — needs COCO per-keypoint sigmas and segmentation-area scale (undefined for MoCap GT)", {"size": 15}),
    ("Pixel MPJPE — not comparable across subjects, distances, resolutions", {"size": 15}),
    ("NMPJPE keeps the full error distribution, is scale-free, and reads as a fraction of body size", {"size": 15, "bold": True}),
], size=16)

# B3 Rotation
s = backup("B3")
title(s, "Backup: viewpoint annotation and calibration", ACCENT)
text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.3), [
    ("θ_raw = atan2(|Δz|, |Δx|) over the 3D ground-truth shoulder vector — folded to [0°, 90°]", {"size": 17, "bold": True}),
    ("From ground truth, not predictions → no circularity (model errors cannot distort the angle bins)", {"size": 16}),
    ("Camera-relative offset: θ_c17 = |θ_raw − 65°|, θ_c18 = 90° − θ_c17", {"size": 16}),
    ("65° calibrated empirically from 3 frontal reference recordings (PM_114, PM_122, PM_109); estimated uncertainty ±5° — REHAB24-6 publishes no full extrinsics", {"size": 16}),
    ("10° bins are wider than the ±5° uncertainty; the offset is global → any error shifts all models identically, the between-model comparison is invariant", {"size": 16, "bold": True}),
    ("Viewpoint distribution is strongly bimodal (40.8% frontal / 17.5% intermediate / 41.6% lateral) → interpreted as viewpoint sensitivity under fixed camera geometry, not a dense rotation study", {"size": 16}),
], size=16)

# B4 Mapping
s = backup("B4")
title(s, "Backup: keypoint mapping (33 / 17 / 26 → 12)", ACCENT)
rows_map = [
    ("Joint", "GT joint", "COCO idx", "MediaPipe idx"),
    ("Shoulders L/R", "LeftArm / RightArm", "5 / 6", "11 / 12"),
    ("Elbows L/R", "LeftForeArm / RightForeArm", "7 / 8", "13 / 14"),
    ("Wrists L/R", "LeftHand / RightHand", "9 / 10", "15 / 16"),
    ("Hips L/R", "LeftUpLeg / RightUpLeg", "11 / 12", "23 / 24"),
    ("Knees L/R", "LeftLeg / RightLeg", "13 / 14", "25 / 26"),
    ("Ankles L/R", "LeftFoot / RightFoot", "15 / 16", "27 / 28"),
]
y = Inches(1.5)
for i, (a, b, c, d) in enumerate(rows_map):
    bold = i == 0
    fill = RGBColor(0xE3, 0xEB, 0xF7) if i == 0 else (LIGHT if i % 2 else RGBColor(0xFF, 0xFF, 0xFF))
    for j, (val, w) in enumerate(zip((a, b, c, d), (Inches(2.6), Inches(4.6), Inches(2.0), Inches(2.4)))):
        x = Inches(0.6) + sum([Inches(2.6), Inches(4.6), Inches(2.0)][:j], Emu(0))
        box(s, x, y, w, Inches(0.55), fill, None, [(val, {"size": 13, "bold": bold})], align=PP_ALIGN.LEFT)
    y += Inches(0.55)
text(s, Inches(0.5), Inches(5.75), Inches(12.3), Inches(1.2), [
    ("MoCap naming convention: a joint carries the name of the segment that STARTS at it — LeftArm = shoulder, LeftFoot = ankle", {"size": 15, "bold": True}),
    ("Facial landmarks excluded (no GT markers). Mapping validated visually by overlaying GT and predictions.", {"size": 15}),
], size=15)

# B5 Varianten
s = backup("B5")
title(s, "Backup: all nine variants (accuracy + CPU speed)", ACCENT)
var_rows = [
    ("Variant", "Mean NMPJPE", "FPS (CPU)"),
    ("MoveNet MultiPose", "10.52%", "27.7"),
    ("YOLOv8 Medium", "10.86%", "4.6"),
    ("YOLOv8 Small", "11.43%", "10.3"),
    ("MediaPipe Heavy", "11.51%", "6.0"),
    ("MediaPipe Full", "12.53%", "14.7"),
    ("YOLOv8 Nano", "12.77%", "18.9"),
    ("MoveNet SP Thunder", "13.33%", "27.6"),
    ("MediaPipe Lite", "13.54%", "21.9"),
    ("MoveNet SP Lightning", "14.97%", "100.6"),
]
y = Inches(1.45)
for i, (a, b, c) in enumerate(var_rows):
    bold = i == 0
    fill = RGBColor(0xE3, 0xEB, 0xF7) if i == 0 else (LIGHT if i % 2 else RGBColor(0xFF, 0xFF, 0xFF))
    box(s, Inches(1.2), y, Inches(4.6), Inches(0.5), fill, None, [(a, {"size": 13, "bold": bold})], align=PP_ALIGN.LEFT)
    box(s, Inches(5.8), y, Inches(2.6), Inches(0.5), fill, None, [(b, {"size": 13, "bold": bold})])
    box(s, Inches(8.4), y, Inches(2.6), Inches(0.5), fill, None, [(c, {"size": 13, "bold": bold})])
    y += Inches(0.5)
text(s, Inches(1.2), Inches(6.6), Inches(11.0), Inches(0.5),
     [("MultiPose > SP variants: SP models expect a centered, frame-filling person — without intelligent cropping the subject is small in the input", {"size": 14, "color": GRAY})])

# B6 COCO
s = backup("B6")
title(s, "Backup: COCO auxiliary protocol", ACCENT)
text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.3), [
    ("1,519 val2017 images with ≥1 annotated person having ≥6 labeled keypoints", {"size": 16}),
    ("Predicted instance matched to the GT person with the nearest torso center (works for all models — MediaPipe has no boxes; OKS matching would be circular)", {"size": 16}),
    ("Same 12 evaluated joints; only visible GT joints included", {"size": 16}),
    ("On REHAB24-6, no GT matching is used — the pipeline selects the most prominent person, as a deployed app would (deployment realism)", {"size": 16, "bold": True}),
    ("Because matching and visibility handling differ, COCO results are descriptive: a transferability check under a different protocol — deliberately no inferential claims", {"size": 16}),
], size=16)

# B7 Coach per video
s = backup("B7")
title(s, "Backup: coach subset per video (mean NMPJPE)", ACCENT)
coach_rows = [
    ("Sequence", "MediaPipe", "MoveNet", "YOLOv8"),
    ("PM_010-c17", "81.9%", "70.7%", "69.9%"),
    ("PM_011-c17", "38.0%", "59.7%", "66.2%"),
    ("PM_108-c17", "17.2%", "46.2%", "51.1%"),
    ("PM_119-c17", "24.1%", "74.0%", "73.0%"),
    ("PM_121-c17", "30.9%", "60.3%", "75.3%"),
]
y = Inches(1.5)
for i, (a, b, c, d) in enumerate(coach_rows):
    bold = i == 0
    fill = RGBColor(0xE3, 0xEB, 0xF7) if i == 0 else (LIGHT if i % 2 else RGBColor(0xFF, 0xFF, 0xFF))
    box(s, Inches(1.6), y, Inches(2.8), Inches(0.62), fill, None, [(a, {"size": 14, "bold": bold})], align=PP_ALIGN.LEFT)
    box(s, Inches(4.4), y, Inches(2.4), Inches(0.62), fill, None, [(b, {"size": 14, "bold": bold, "color": MP_GREEN if not bold else DARK})])
    box(s, Inches(6.8), y, Inches(2.4), Inches(0.62), fill, None, [(c, {"size": 14, "bold": bold, "color": MN_RED if not bold else DARK})])
    box(s, Inches(9.2), y, Inches(2.4), Inches(0.62), fill, None, [(d, {"size": 14, "bold": bold, "color": YL_BLUE if not bold else DARK})])
    y += Inches(0.62)
text(s, Inches(1.6), Inches(5.6), Inches(10.5), Inches(1.2), [
    ("PM_010 breaks all three models (>65%) — no selection strategy is perfect", {"size": 15, "bold": True}),
    ("MediaPipe stays below 40% in the other four sequences. NMPJPE before the outlier cut (failure regime). Named limitation: n=5, descriptive only.", {"size": 14, "color": GRAY}),
], size=15)

# B8 Hips
s = backup("B8")
title(s, "Backup: per-joint errors and the hip effect", ACCENT)
pj_rows = [
    ("Region", "MediaPipe", "MoveNet", "YOLOv8"),
    ("Shoulders", "8.76%", "10.29%", "10.21%"),
    ("Elbows", "11.24%", "10.00%", "13.08%"),
    ("Wrists", "12.01%", "11.84%", "16.25%"),
    ("Hips", "16.72%", "11.28%", "14.59%"),
    ("Knees", "9.92%", "7.76%", "9.87%"),
    ("Ankles", "14.59%", "11.41%", "12.59%"),
]
y = Inches(1.45)
for i, (a, b, c, d) in enumerate(pj_rows):
    bold = i == 0
    hip = a == "Hips"
    fill = RGBColor(0xE3, 0xEB, 0xF7) if i == 0 else (RGBColor(0xFB, 0xE9, 0xE7) if hip else (LIGHT if i % 2 else RGBColor(0xFF, 0xFF, 0xFF)))
    box(s, Inches(1.6), y, Inches(2.8), Inches(0.55), fill, None, [(a, {"size": 13, "bold": bold or hip})], align=PP_ALIGN.LEFT)
    box(s, Inches(4.4), y, Inches(2.4), Inches(0.55), fill, None, [(b, {"size": 13, "bold": bold or hip})])
    box(s, Inches(6.8), y, Inches(2.4), Inches(0.55), fill, None, [(c, {"size": 13, "bold": bold})])
    box(s, Inches(9.2), y, Inches(2.4), Inches(0.55), fill, None, [(d, {"size": 13, "bold": bold})])
    y += Inches(0.55)
text(s, Inches(1.6), Inches(5.35), Inches(10.5), Inches(1.5), [
    ("Excluding hips: MediaPipe 12.53% → 11.66%; gap to MoveNet shrinks 2.01pp → 1.21pp", {"size": 15, "bold": True}),
    ("Error concentrated in one structurally specific region → consistent with a landmark-definition mismatch (BlazePose body model vs. MoCap skeleton) — indirect evidence; the index mapping itself is validated", {"size": 14, "color": GRAY}),
], size=15)

# B9 Architektur
s = backup("B9")
title(s, "Backup: model internals", ACCENT)
text(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(5.5), [
    ("BlazePose (MediaPipe)", {"bold": True, "size": 17, "color": MP_GREEN}),
    ("Two-stage detector–tracker; the paper describes the person detector as face-based (BlazeFace-derived) with alignment keypoints (mid-hip, scale, rotation). Training combines heatmap + regression supervision; inference is regression-only. 33 landmarks. IMAGE mode = detector runs every frame.", {"size": 14}),
    ("MoveNet", {"bold": True, "size": 17, "color": MN_RED}),
    ("MobileNetV2 + FPN backbone, CenterNet-style: 4 heads — person-center heatmap, keypoint regression field, keypoint heatmaps, local offsets. Decoding weights heatmap maxima by inverse distance to the regressed position. Lightning 192², Thunder 256², MultiPose up to 6 persons + boxes. Trained on COCO + Google-internal 'Active' set (yoga/fitness/dance — per TF blog).", {"size": 14}),
    ("YOLOv8-Pose", {"bold": True, "size": 17, "color": YL_BLUE}),
    ("Anchor-free one-stage detector (C2f blocks, decoupled head) extended with a pose head: per detection box + 17 keypoints (x, y, conf) in one pass; OKS-based keypoint loss. Ultralytics PyTorch runtime.", {"size": 14}),
], size=15)

# ---- speichern ----
out = os.path.join(DEF, "defense_slides.pptx")
prs.save(out)
print("Gespeichert:", out)
print("Folien gesamt:", len(prs.slides.__iter__.__self__._sldIdLst))
